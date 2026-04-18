"""
Vatic Verticalization PRD — PowerPoint Generator
Generates a dark, editorial-luxury branded presentation from VERTICALIZATION_PRD.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ─── Brand Colors ────────────────────────────────────────────────────────────
BG         = RGBColor(0x09, 0x09, 0x0B)   # zinc-950 — slide background
CARD       = RGBColor(0x18, 0x18, 0x1B)   # zinc-900 — card/panel fills
BORDER     = RGBColor(0x3F, 0x3F, 0x46)   # zinc-600 — borders/dividers
WHITE      = RGBColor(0xFA, 0xFA, 0xFA)   # near-white — primary text
MUTED      = RGBColor(0xA1, 0xA1, 0xAA)   # zinc-400  — secondary text
VIOLET     = RGBColor(0x7C, 0x3A, 0xED)   # violet-600 — Vatic primary accent
VIOLET_LT  = RGBColor(0xA7, 0x8B, 0xFA)   # violet-400 — lighter accent
CYAN       = RGBColor(0x06, 0xB6, 0xD4)   # cyan-500   — gradient complement
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)   # amber-500  — warning / important
GREEN      = RGBColor(0x10, 0xB9, 0x81)   # emerald-500 — "no change needed"
RED        = RGBColor(0xEF, 0x44, 0x44)   # red-500    — critical

# ─── Dimensions (16:9, inches) ───────────────────────────────────────────────
W = 13.33
H = 7.5

# ─── Helpers ─────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # completely blank
    slide = prs.slides.add_slide(blank_layout)
    # dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def add_rect(slide, x, y, w, h, fill_color=None, border_color=None, border_width=None, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.fill.background()  # no line by default
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 0.75)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=16, bold=False, color=None, align=PP_ALIGN.LEFT,
             italic=False, wrap=True, valign=None):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or WHITE
    run.font.name = "Calibri"
    return txBox


def add_text_multi(slide, lines, x, y, w, h,
                   default_size=14, default_color=None, default_bold=False,
                   line_spacing=None, align=PP_ALIGN.LEFT):
    """
    lines: list of dicts with keys: text, size (opt), bold (opt), color (opt), italic (opt)
    """
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.space_after = Pt(line_spacing)
        run = p.add_run()
        run.text = line.get("text", "")
        run.font.size = Pt(line.get("size", default_size))
        run.font.bold = line.get("bold", default_bold)
        run.font.italic = line.get("italic", False)
        run.font.color.rgb = line.get("color", default_color or WHITE)
        run.font.name = "Calibri"
    return txBox


def add_accent_bar(slide, y=0.55, height=0.04):
    """Add the thin violet→cyan gradient accent bar at the top of content slides."""
    # Simulate gradient with two overlapping rects (violet → cyan)
    bar1 = add_rect(slide, 0, y, W * 0.5, height, fill_color=VIOLET)
    bar2 = add_rect(slide, W * 0.5, y, W * 0.5, height, fill_color=CYAN)
    return bar1, bar2


def slide_header(slide, title, subtitle=None, title_y=0.18, bar_y=0.55):
    add_accent_bar(slide, y=bar_y)
    add_text(slide, title, 0.5, title_y, W - 1, 0.55,
             size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.5, title_y + 0.42, W - 1, 0.35,
                 size=13, color=MUTED)


def card(slide, x, y, w, h, title=None, title_color=None, body_lines=None, accent_color=None):
    """Draw a dark card with optional top-accent stripe, title, and body lines."""
    add_rect(slide, x, y, w, h, fill_color=CARD, border_color=BORDER, border_width=0.5)
    if accent_color:
        add_rect(slide, x, y, w, 0.06, fill_color=accent_color)
    if title:
        add_text(slide, title, x + 0.15, y + 0.12, w - 0.3, 0.35,
                 size=12, bold=True, color=title_color or WHITE)
    if body_lines:
        top_offset = 0.12 + (0.3 if title else 0)
        add_text_multi(slide, body_lines, x + 0.15, y + top_offset, w - 0.3,
                       h - top_offset - 0.1, default_size=10, default_color=MUTED)


def badge(slide, text, x, y, w=1.0, h=0.28, bg=VIOLET, text_color=WHITE, size=9):
    add_rect(slide, x, y, w, h, fill_color=bg)
    add_text(slide, text, x, y + 0.03, w, h - 0.05,
             size=size, bold=True, color=text_color, align=PP_ALIGN.CENTER)


def bullet_list(slide, items, x, y, w, item_size=12, color=None, bullet_color=None, spacing=0.32):
    bc = bullet_color or VIOLET_LT
    tc = color or WHITE
    for i, item in enumerate(items):
        # bullet dot
        add_text(slide, "▸", x, y + i * spacing, 0.2, 0.28, size=item_size, color=bc, bold=True)
        add_text(slide, item, x + 0.22, y + i * spacing, w - 0.22, 0.28,
                 size=item_size, color=tc)


def two_col_bullets(slide, left_items, right_items, x=0.5, y=1.15, col_w=5.9,
                    gap=0.4, item_size=11, spacing=0.3):
    bullet_list(slide, left_items,  x,                   y, col_w, item_size=item_size, spacing=spacing)
    bullet_list(slide, right_items, x + col_w + gap,     y, col_w, item_size=item_size, spacing=spacing)


# ─── Slide Builders ──────────────────────────────────────────────────────────

def slide_title(prs):
    slide = blank_slide(prs)

    # Large gradient accent stripe left edge
    add_rect(slide, 0, 0, 0.06, H, fill_color=VIOLET)

    # Faint right-side glow rectangle
    add_rect(slide, W - 3.5, 1.5, 3.5, 4.5, fill_color=RGBColor(0x12, 0x0C, 0x2A))

    # VATIC wordmark
    add_text(slide, "VATIC", 0.5, 1.6, 4, 0.8,
             size=48, bold=True, color=VIOLET, align=PP_ALIGN.LEFT)

    add_text(slide, "Competitive Intelligence Platform", 0.5, 2.35, 8, 0.4,
             size=15, color=MUTED, italic=True)

    # Divider
    add_rect(slide, 0.5, 2.9, 8, 0.025, fill_color=BORDER)

    add_text(slide, "Verticalization PRD", 0.5, 3.05, 10, 0.7,
             size=36, bold=True, color=WHITE)

    add_text(slide, "Codebase Audit · Architecture Options · Recommendation", 0.5, 3.72, 10, 0.4,
             size=14, color=VIOLET_LT, italic=True)

    add_text(slide, "April 7, 2026  ·  Anand Iyer  ·  Alive Labs / Prophet", 0.5, H - 0.7, 10, 0.3,
             size=10, color=MUTED)

    # Decorative corner badge
    badge(slide, "RESEARCH DRAFT", W - 2.2, H - 0.55, w=1.9, h=0.28,
          bg=RGBColor(0x3F, 0x3F, 0x46), text_color=MUTED, size=8)


def slide_context(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Background & Strategic Context",
                 "What led to this PRD and what decisions were already made")

    # Two cards
    card(slide, 0.4, 0.75, 5.9, 2.6,
         title="The April 6 Meeting Decisions",
         title_color=VIOLET_LT,
         accent_color=VIOLET,
         body_lines=[
             {"text": "1.  Pivot primary target vertical from restaurants → liquor stores",
              "size": 11, "color": WHITE},
             {"text": "    Hunter validated concept; liquor store owners constantly analyze competitor pricing",
              "size": 10, "color": MUTED},
             {"text": "2.  Single codebase + shared database architecture",
              "size": 11, "color": WHITE, "bold": True},
             {"text": "    Add industry_type column · branded subdomains per vertical",
              "size": 10, "color": MUTED},
             {"text": "3.  Henry + Anand to present options at next week's meeting",
              "size": 11, "color": AMBER},
         ])

    card(slide, 6.6, 0.75, 6.3, 2.6,
         title="The Two-Sided Revenue Opportunity",
         title_color=CYAN,
         accent_color=CYAN,
         body_lines=[
             {"text": "Side 1 — Subscriptions", "size": 11, "bold": True, "color": WHITE},
             {"text": "Competitive intelligence SaaS for liquor store owners",
              "size": 10, "color": MUTED},
             {"text": "Side 2 — Brand Intelligence (new)", "size": 11, "bold": True, "color": WHITE},
             {"text": "Aggregate competitor pricing + sales data → sell back to",
              "size": 10, "color": MUTED},
             {"text": "distributors (Bacardi, etc.) at data scale",
              "size": 10, "color": MUTED},
             {"text": "This second play does not exist in the restaurant sector",
              "size": 10, "color": AMBER, "italic": True},
         ])

    # Current state bar
    add_rect(slide, 0.4, 3.55, 12.5, 0.025, fill_color=BORDER)
    add_text(slide, "Current State  —  10 Phases Shipped", 0.4, 3.65, 5, 0.3,
             size=11, bold=True, color=MUTED)

    phases = [
        "Auth & Onboarding", "Location Mgmt", "Competitor Discovery",
        "Snapshot Pipeline", "Insight Engine", "Billing / Stripe",
        "Background Jobs", "Social Intelligence", "Social Visual Intel",
        "Actionable Insights",
    ]
    chip_w = 1.95
    chip_gap = 0.08
    per_row = 5
    for i, p in enumerate(phases):
        col = i % per_row
        row = i // per_row
        cx = 0.4 + col * (chip_w + chip_gap)
        cy = 4.05 + row * 0.42
        add_rect(slide, cx, cy, chip_w, 0.32, fill_color=RGBColor(0x1E, 0x1A, 0x30),
                 border_color=VIOLET, border_width=0.5)
        add_text(slide, f"✓  {p}", cx + 0.1, cy + 0.04, chip_w - 0.15, 0.25,
                 size=9, color=VIOLET_LT)

    # NOT shipped
    add_text(slide, "Not yet shipped:", 0.4, 6.95, 3, 0.3, size=10, bold=True, color=MUTED)
    not_shipped = ["Ask Prophet (LLM Chat)", "Data Retention", "Team Management"]
    for i, item in enumerate(not_shipped):
        add_text(slide, f"○  {item}", 3.2 + i * 3.2, 6.95, 3, 0.3, size=10, color=RED)


def slide_guiding_principles(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Guiding Principles",
                 "Constraints that shape what good verticalization architecture looks like")

    principles = [
        ("Speed to Market",      "Liquor stores are ready now. Architecture must not block a fast first deployment.",   VIOLET),
        ("Maintainability",      "One engineering team. Diverging codebases create exponential maintenance overhead.",   CYAN),
        ("Data Separation",      "If a vertical is sold, buyer expects clean data. Don't make separation impossible.",   AMBER),
        ("Performance",          "Shared database must not create query issues as verticals scale.",                     GREEN),
        ("Brand Differentiation","Each vertical needs distinct landing pages, copy, and onboarding.",                   VIOLET_LT),
        ("API Cost Management",  "Not all signals relevant for all verticals. Signal toggles prevent wasteful API spend.", MUTED),
    ]

    cols = 3
    card_w = (W - 0.5 - 0.4 - (cols - 1) * 0.25) / cols
    card_h = 2.35
    for i, (title, desc, color) in enumerate(principles):
        col = i % cols
        row = i // cols
        cx = 0.4 + col * (card_w + 0.25)
        cy = 0.82 + row * (card_h + 0.25)
        add_rect(slide, cx, cy, card_w, card_h, fill_color=CARD, border_color=BORDER, border_width=0.5)
        add_rect(slide, cx, cy, card_w, 0.07, fill_color=color)
        add_text(slide, title, cx + 0.15, cy + 0.15, card_w - 0.3, 0.38,
                 size=13, bold=True, color=WHITE)
        add_text(slide, desc, cx + 0.15, cy + 0.55, card_w - 0.3, card_h - 0.65,
                 size=10, color=MUTED, wrap=True)


def slide_audit_overview(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Codebase Audit — What We Analyzed",
                 "10 dimensions audited across the full Prophet codebase")

    rows = [
        ("Database Schema",          "Largely industry-agnostic",               GREEN,  "LOW"),
        ("TypeScript Types",          "Most hardcoded layer — MenuType, DetectedFeatures, Cuisine", RED, "CRITICAL"),
        ("Intelligence Signals",      "7 of 8 signals are fully generic. Only Content/Menu is specific.", AMBER, "MEDIUM"),
        ("AI / Gemini Prompts",       "1 critical prompt (menu extraction) fully hardcoded. Others generic.", AMBER, "MEDIUM"),
        ("Server Actions / API",      "3 actions have meaningful hardcoding",    AMBER,  "MEDIUM"),
        ("Insight Rules Engine",      "8 content rules restaurant-specific. ~40+ others are generic.", RED, "HIGH"),
        ("Onboarding Flow",           "Most concentrated user-facing restaurant language in the app.", RED, "HIGH"),
        ("UI / Landing Page",         "Dashboard is mostly generic. Landing page is restaurant-specific.", AMBER, "MEDIUM"),
        ("Email Templates",           "One welcome template has restaurant-specific copy. Others generic.", GREEN, "LOW"),
        ("Config & Constants",        "Menu patterns, promo keywords, cuisine list all hardcoded.", RED, "HIGH"),
    ]

    row_h = 0.49
    col_x = [0.35, 4.2, 9.5, 11.2]
    col_w = [3.7, 5.1, 1.5, 1.9]

    # Header row
    headers = ["Category", "Verdict", "Severity", "Blocks Launch?"]
    for j, hdr in enumerate(headers):
        add_rect(slide, col_x[j], 0.78, col_w[j], 0.35, fill_color=RGBColor(0x27, 0x27, 0x2A))
        add_text(slide, hdr, col_x[j] + 0.1, 0.82, col_w[j] - 0.1, 0.28,
                 size=9, bold=True, color=MUTED)

    severity_bg = {
        "CRITICAL": RED,
        "HIGH":     AMBER,
        "MEDIUM":   RGBColor(0x0E, 0x70, 0x90),
        "LOW":      RGBColor(0x06, 0x5F, 0x46),
    }

    for i, (cat, verdict, color, severity) in enumerate(rows):
        ry = 1.17 + i * row_h
        row_bg = CARD if i % 2 == 0 else RGBColor(0x1C, 0x1C, 0x1F)
        add_rect(slide, col_x[0], ry, sum(col_w) + 0.2, row_h - 0.03, fill_color=row_bg)
        add_text(slide, cat, col_x[0] + 0.1, ry + 0.1, col_w[0] - 0.15, row_h - 0.15,
                 size=10, bold=True, color=WHITE)
        add_text(slide, verdict, col_x[1] + 0.1, ry + 0.1, col_w[1] - 0.15, row_h - 0.15,
                 size=9.5, color=MUTED)
        sbg = severity_bg.get(severity, BORDER)
        add_rect(slide, col_x[2] + 0.15, ry + 0.1, 1.2, 0.27, fill_color=sbg)
        add_text(slide, severity, col_x[2] + 0.15, ry + 0.11, 1.2, 0.26,
                 size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        blocks = "YES" if severity in ("CRITICAL", "HIGH") else "No"
        bc = RED if blocks == "YES" else GREEN
        add_text(slide, blocks, col_x[3] + 0.3, ry + 0.1, 1.3, 0.26,
                 size=10, bold=True, color=bc, align=PP_ALIGN.CENTER)


def slide_db_schema(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Audit: Database Schema",
                 "The best news — the schema is almost entirely vertical-agnostic")

    badge(slide, "LOW RISK", W - 2.0, 0.18, w=1.7, h=0.3, bg=GREEN, text_color=WHITE, size=9)

    add_text(slide, "Generic Tables (No Changes Required)", 0.4, 0.78, 7, 0.3,
             size=12, bold=True, color=GREEN)
    generic_tables = [
        "organizations  —  name, tier, trial dates, stripe_customer_id",
        "locations  —  name, address, google_place_id, website",
        "competitors  —  name, place_id, status (approved / ignored)",
        "snapshots  —  polymorphic JSON blobs keyed by snapshot_type",
        "insights / social_profiles / seo / events / photos / busy_times / weather",
    ]
    bullet_list(slide, generic_tables, 0.4, 1.08, 8.5, item_size=10.5, spacing=0.31)

    add_rect(slide, 0.4, 2.73, 12.5, 0.025, fill_color=BORDER)

    add_text(slide, "Critical Gap  —  Missing industry_type Column", 0.4, 2.82, 9, 0.3,
             size=12, bold=True, color=RED)

    # Three action cards
    actions = [
        ("Add industry_type column",
         "organizations table\ncheck constraint: ('restaurant', 'liquor_store')\nIndex for query performance",
         VIOLET),
        ("Backfill existing orgs",
         "Set industry_type = 'restaurant'\nfor all existing organizations\nLow-risk additive migration",
         AMBER),
        ("Optional: vertical_config JSONB",
         "Per-org signal enable/disable overrides\nStores per-org custom config\nAdditive, non-breaking",
         CYAN),
    ]
    card_w = 3.9
    for i, (title, body, color) in enumerate(actions):
        cx = 0.4 + i * (card_w + 0.27)
        add_rect(slide, cx, 3.15, card_w, 3.7, fill_color=CARD, border_color=color, border_width=1.0)
        add_rect(slide, cx, 3.15, card_w, 0.07, fill_color=color)
        add_text(slide, title, cx + 0.15, 3.25, card_w - 0.3, 0.38,
                 size=11, bold=True, color=WHITE)
        add_text(slide, body, cx + 0.15, 3.65, card_w - 0.3, 2.8,
                 size=10, color=MUTED, wrap=True)

    add_text(slide, "Migration complexity: LOW  —  These are additive changes. No existing data at risk.",
             0.4, 7.1, 12, 0.3, size=10, bold=False, color=GREEN, italic=True)


def slide_type_system(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Audit: TypeScript Type System",
                 "The most hardcoded layer — requires the most thoughtful refactoring")

    badge(slide, "CRITICAL", W - 2.0, 0.18, w=1.7, h=0.3, bg=RED, size=9)

    # Left: current hardcoded types
    add_text(slide, "Current Hardcoded Restaurant Types", 0.4, 0.75, 6.2, 0.3,
             size=12, bold=True, color=RED)

    type_items = [
        ('MenuType', '"dine_in" | "catering" | "banquet" | "happy_hour" | "kids" | "other"', RED),
        ('MenuItem.tags', '"vegan" | "spicy" | "gluten-free"  — food-specific tags', RED),
        ('DetectedFeatures', 'reservation, privateDining, catering, happyHour, deliveryPlatforms', RED),
        ('CorePage.type', '"reservations" | "catering" | "menu" | "about" | …', AMBER),
        ('CUISINES[]', '"American", "Italian", "Mexican", "Asian", "Bar & Grill", …', RED),
    ]
    for i, (name, desc, color) in enumerate(type_items):
        ry = 1.08 + i * 0.52
        add_rect(slide, 0.4, ry, 6.1, 0.46, fill_color=CARD, border_color=color, border_width=0.75)
        add_text(slide, name, 0.6, ry + 0.06, 1.8, 0.32,
                 size=10, bold=True, color=color)
        add_text(slide, desc, 2.45, ry + 0.06, 3.9, 0.32,
                 size=9.5, color=MUTED)

    # Right: target vertical-agnostic types
    add_text(slide, "Target Vertical-Configurable Types", 6.9, 0.75, 6.0, 0.3,
             size=12, bold=True, color=GREEN)

    target_items = [
        ('ContentCategoryType', 'Loaded from VerticalConfig — not a hardcoded union', GREEN),
        ('CatalogItem.tags',    'Vertical-specific attribute array (e.g. "ABV" for liquor)', GREEN),
        ('DetectedFeatures',    'Polymorphic interface per vertical (FeatureDefinition[])', GREEN),
        ('CorePage.type',       'Configurable set per vertical via VerticalConfig', GREEN),
        ('businessCategories[]','Vertical-specific list (cuisines OR store types)', GREEN),
    ]
    for i, (name, desc, color) in enumerate(target_items):
        ry = 1.08 + i * 0.52
        add_rect(slide, 6.9, ry, 6.0, 0.46, fill_color=CARD, border_color=color, border_width=0.75)
        add_text(slide, name, 7.1, ry + 0.06, 2.0, 0.32,
                 size=10, bold=True, color=color)
        add_text(slide, desc, 9.15, ry + 0.06, 3.6, 0.32,
                 size=9.5, color=MUTED)

    # Arrow in middle
    add_text(slide, "→", 6.4, 2.8, 0.5, 0.5, size=24, bold=True, color=VIOLET, align=PP_ALIGN.CENTER)

    add_rect(slide, 0.4, 3.8, 12.5, 0.025, fill_color=BORDER)
    add_text(slide,
             "Key insight: No type change requires a database migration — all changes are purely in TypeScript and application logic.",
             0.4, 3.88, 12.2, 0.3, size=10, color=MUTED, italic=True)


def slide_signals(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Audit: Intelligence Signals",
                 "7 of 8 signals require zero changes — only the Content/Menu signal is vertical-specific")

    # Generic signals grid
    add_text(slide, "Generic Signals — No Changes Required  (7/8)", 0.4, 0.75, 8, 0.3,
             size=12, bold=True, color=GREEN)

    generic_signals = [
        ("Competitor Monitoring", "Google Places\nReviews, ratings, hours, attributes"),
        ("SEO / Search Visibility", "DataForSEO\n12 endpoints, keyword rankings"),
        ("Local Events", "DataForSEO Events SERP\nNearby event matching"),
        ("Foot Traffic", "Outscraper\nPopular Times, hourly patterns"),
        ("Weather Intelligence", "OpenWeatherMap\nHistorical + forecast"),
        ("Social Media", "Data365\nInstagram, Facebook, TikTok"),
        ("Photo Intelligence", "Google Places + Gemini Vision\nAmbiance, quality analysis"),
    ]
    sg_w = 1.82
    sg_h = 1.0
    gap = 0.08
    for i, (title, desc) in enumerate(generic_signals):
        cx = 0.4 + i * (sg_w + gap)
        add_rect(slide, cx, 1.1, sg_w, sg_h, fill_color=CARD, border_color=GREEN, border_width=0.5)
        add_rect(slide, cx, 1.1, sg_w, 0.05, fill_color=GREEN)
        add_text(slide, title, cx + 0.1, 1.17, sg_w - 0.2, 0.3, size=9, bold=True, color=WHITE)
        add_text(slide, desc, cx + 0.1, 1.48, sg_w - 0.2, 0.58, size=8, color=MUTED)

    add_rect(slide, 0.4, 2.22, 12.5, 0.025, fill_color=BORDER)

    # The problem signal
    add_text(slide, "Restaurant-Specific Signal — Full Vertical Abstraction Required  (1/8)",
             0.4, 2.3, 12, 0.3, size=12, bold=True, color=RED)

    # Two columns: restaurant vs liquor store
    add_rect(slide, 0.4, 2.65, 5.9, 4.5, fill_color=CARD, border_color=RED, border_width=1.0)
    add_rect(slide, 0.4, 2.65, 5.9, 0.07, fill_color=RED)
    add_text(slide, "Content / Menu Signal  (Restaurant)", 0.55, 2.72, 5.5, 0.3,
             size=11, bold=True, color=WHITE)
    restaurant_items = [
        "Scrapes website for a menu page",
        "Extracts menu items: name, price, tags (vegan, spicy)",
        "Classifies into MenuType: dine_in, catering, happy_hour, kids",
        "Detects: reservations, online ordering, private dining",
        "Detects: DoorDash / Grubhub / UberEats integrations",
        "Generates 8 menu-specific insight rules",
    ]
    bullet_list(slide, restaurant_items, 0.55, 3.08, 5.6, item_size=10, spacing=0.3)

    add_text(slide, "→", 6.4, 4.6, 0.5, 0.5, size=24, bold=True, color=VIOLET, align=PP_ALIGN.CENTER)

    add_rect(slide, 6.9, 2.65, 5.9, 4.5, fill_color=CARD, border_color=CYAN, border_width=1.0)
    add_rect(slide, 6.9, 2.65, 5.9, 0.07, fill_color=CYAN)
    add_text(slide, "Content / Catalog Signal  (Liquor Store)", 7.05, 2.72, 5.5, 0.3,
             size=11, bold=True, color=WHITE)
    liquor_items = [
        "Scrapes website for a products / spirits page",
        "Extracts: spirit type, brand, price, ABV, size, country",
        "Classifies into: spirits, wine, beer, mixers, accessories",
        "Detects: curbside pickup, delivery, tasting events",
        "Detects: Drizly / Instacart / GoPuff integrations",
        "Generates 6 catalog-specific insight rules",
    ]
    bullet_list(slide, liquor_items, 7.05, 3.08, 5.6, item_size=10,
                color=WHITE, bullet_color=CYAN, spacing=0.3)


def slide_prompts(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Audit: AI / Gemini Prompts",
                 "One critical prompt is deeply restaurant-specific — the rest are already generic or need minor injection")

    rows = [
        ("Menu Extraction",         "lib/ai/gemini.ts",
         "CRITICAL — fully restaurant-specific",
         '"restaurant menu", "appetizers", "entrees", "catering", "happy_hour", "kids menus"',
         "Full rewrite per vertical", RED),
        ("Competitor Discovery",    "lib/providers/gemini.ts",
         "Generic — already uses 'local business'",
         "No restaurant-specific nouns. Uses generic business language throughout.",
         "None required", GREEN),
        ("Insight Narratives",      "lib/ai/prompts/insights.ts",
         "Generic — 'local businesses'",
         "Already uses generic framing. Add vertical context variable for richer output.",
         "Inject industry_type variable", AMBER),
        ("Priority Briefing",       "lib/ai/prompts/priority-briefing.ts",
         "Generic — 'local businesses'",
         "Generic framing. Adding 'local liquor store' context would improve output quality.",
         "Inject industry_type variable", AMBER),
        ("Prophet Chat",            "lib/ai/prompts/prophet-chat.ts",
         "Not yet active",
         "Chat endpoint is scaffolded but does not call any LLM. Design with vertical context from start.",
         "Design with vertical context", CYAN),
    ]

    row_h = 1.12
    for i, (name, file, status, detail, action, color) in enumerate(rows):
        ry = 0.78 + i * row_h
        add_rect(slide, 0.35, ry, 12.6, row_h - 0.08, fill_color=CARD, border_color=BORDER)
        add_rect(slide, 0.35, ry, 0.07, row_h - 0.08, fill_color=color)
        add_text(slide, name, 0.55, ry + 0.08, 2.2, 0.3, size=11, bold=True, color=WHITE)
        add_text(slide, file, 0.55, ry + 0.4, 2.4, 0.25, size=8, color=MUTED, italic=True)
        add_text(slide, status, 3.0, ry + 0.08, 3.8, 0.3, size=10, bold=True, color=color)
        add_text(slide, detail, 3.0, ry + 0.42, 5.8, 0.55, size=9, color=MUTED)
        add_text(slide, action, 9.1, ry + 0.22, 3.6, 0.3, size=10, bold=True, color=color)


def slide_onboarding(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Audit: Onboarding Flow",
                 "The most user-facing concentration of restaurant language — mechanical to fix, but touches many files")

    badge(slide, "HIGH", W - 1.8, 0.18, w=1.5, h=0.3, bg=AMBER, size=9)

    steps = [
        ("Step 0\nSplash", '"Set up my restaurant"', RED, [
            '"Set up my restaurant"  →  "Set up my business"',
        ]),
        ("Step 1\nBusiness Info", "restaurantName, CUISINES[]", RED, [
            'restaurantName  →  businessName',
            '"Cuisine Type"  →  "Business Type"',
            'CUISINES[]  →  verticalConfig.businessCategories',
            '"Your Restaurant"  →  verticalConfig.businessLabel',
        ]),
        ("Step 2\nCompetitor Select", '"nearby restaurants"', RED, [
            '"Searching for nearby restaurants…"',
            '"We found nearby restaurants. Pick up to 5…"',
            'Food emojis  →  verticalConfig.categoryEmojis',
        ]),
        ("Step 3\nIntel Settings", "One restaurant reference", AMBER, [
            '"Get alerted when new restaurants open…"',
            '→ Use verticalConfig.competitorLabel',
        ]),
    ]

    step_w = 3.05
    step_gap = 0.18
    for i, (step_title, subtitle, color, changes) in enumerate(steps):
        cx = 0.35 + i * (step_w + step_gap)
        # Step card
        add_rect(slide, cx, 0.75, step_w, 1.2, fill_color=CARD, border_color=color, border_width=0.75)
        add_rect(slide, cx, 0.75, step_w, 0.06, fill_color=color)
        add_text(slide, step_title, cx + 0.12, 0.83, step_w - 0.25, 0.45,
                 size=10, bold=True, color=WHITE)
        add_text(slide, subtitle, cx + 0.12, 1.28, step_w - 0.25, 0.55,
                 size=9, color=MUTED, italic=True)

        # Changes card below
        add_rect(slide, cx, 2.07, step_w, H - 2.45, fill_color=RGBColor(0x1A, 0x10, 0x10),
                 border_color=color, border_width=0.5)
        add_text(slide, "Changes Required", cx + 0.12, 2.14, step_w - 0.25, 0.26,
                 size=9, bold=True, color=color)
        for j, change in enumerate(changes):
            add_text(slide, f"• {change}", cx + 0.12, 2.42 + j * 0.47, step_w - 0.25, 0.42,
                     size=8.5, color=MUTED)

    # Bottom summary
    add_rect(slide, 0.35, 6.9, 12.6, 0.38, fill_color=RGBColor(0x0F, 0x0F, 0x10))
    add_text(slide,
             "Onboarding State:  restaurantName → businessName   ·   cuisine → businessCategory   ·   "
             "All changes are mechanical renames + config-driven copy. No logic restructuring needed.",
             0.5, 6.95, 12.3, 0.3, size=9.5, color=MUTED)


def slide_insight_rules(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Audit: Insight Rules Engine",
                 "8 content rules are restaurant-specific. ~40+ other rules across 6 signal types are already generic.")

    add_text(slide, "Restaurant-Specific Rules  (Must be rewritten per vertical)",
             0.4, 0.75, 8, 0.28, size=11, bold=True, color=RED)

    rs_rules = [
        ("menu.price_positioning_shift",  "Competitor dropped/raised menu prices"),
        ("menu.category_gap",             "Competitor offers a menu category you don't"),
        ("menu.signature_item_missing",   "Competitor's top item not on your menu"),
        ("menu.promo_signal_detected",    "Promo keywords: 'happy hour', 'kids eat free', 'prix fixe'…"),
        ("menu.menu_change_detected",     "Competitor's menu changed since last snapshot"),
        ("content.conversion_feature_gap","Competitor has: online reservations, private dining, catering"),
        ("content.delivery_platform_gap", "Competitor is on DoorDash/Grubhub/UberEats, you aren't"),
        ("menu.catering_pricing_gap",     "Competitor offers catering packages, you don't"),
    ]
    for i, (rule, desc) in enumerate(rs_rules):
        col = i % 2
        row = i // 2
        rx = 0.4 + col * 6.3
        ry = 1.07 + row * 0.5
        add_rect(slide, rx, ry, 6.1, 0.44, fill_color=RGBColor(0x1A, 0x10, 0x10),
                 border_color=RED, border_width=0.5)
        add_text(slide, rule,  rx + 0.12, ry + 0.05, 2.9, 0.28, size=9, bold=True, color=RED)
        add_text(slide, desc,  rx + 3.05, ry + 0.05, 2.9, 0.28, size=9, color=MUTED)

    add_rect(slide, 0.4, 3.18, 12.5, 0.025, fill_color=BORDER)

    add_text(slide, "Generic Rules — No Changes Required  (~40+ rules)",
             0.4, 3.28, 7, 0.28, size=11, bold=True, color=GREEN)

    generic_groups = [
        ("Competitor", "rating_change, review_velocity, hours_changed, attribute_change"),
        ("SEO (13 types)", "domain_rank, keyword_movement, competitor_overlap, top_pages…"),
        ("Social (15 types)", "engagement_velocity, follower_growth, sentiment_shift…"),
        ("Social Visual (16 types)", "content_strategy, brand_signals, visual_opportunity…"),
        ("Events", "nearby_event_detected, competitor_proximity_event…"),
        ("Traffic", "peak_hour_shift, competitor_busy_comparison…"),
        ("Weather", "weather_suppression, severe_weather_alert…"),
        ("Photos", "photo_quality_gap, ambiance_comparison, food_presentation…"),
    ]

    gg_w = 3.0
    gg_gap = 0.07
    gg_h = 0.85
    per_row = 4
    for i, (name, desc) in enumerate(generic_groups):
        col = i % per_row
        row = i // per_row
        gx = 0.4 + col * (gg_w + gg_gap)
        gy = 3.62 + row * (gg_h + 0.1)
        add_rect(slide, gx, gy, gg_w, gg_h, fill_color=CARD, border_color=GREEN, border_width=0.5)
        add_text(slide, f"✓  {name}", gx + 0.12, gy + 0.08, gg_w - 0.2, 0.28,
                 size=10, bold=True, color=GREEN)
        add_text(slide, desc, gx + 0.12, gy + 0.38, gg_w - 0.2, 0.42,
                 size=8, color=MUTED)


def slide_impact_matrix(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Audit Summary: Impact Matrix",
                 "At a glance — what changes, what stays the same, what blocks the liquor store launch")

    rows_data = [
        ("Database Schema",    "Core tables",          "None",   "None",     GREEN,  "No"),
        ("Database Schema",    "industry_type column", "Low",    "1 migration", AMBER, "YES"),
        ("Type System",        "MenuType, DetectedFeatures", "Medium", "Full abstraction", RED, "YES"),
        ("Type System",        "Cuisine list",         "Low",    "Config-driven", AMBER, "YES"),
        ("Signals",            "7 generic signals",    "None",   "None",     GREEN,  "No"),
        ("Signals",            "Content/Menu signal",  "High",   "Full rewrite", RED, "YES"),
        ("AI Prompts",         "Menu extraction",      "Medium", "Per-vertical prompt", RED, "YES"),
        ("AI Prompts",         "3 generic prompts",    "Low",    "Inject industry_type", GREEN, "No"),
        ("Onboarding",         "Field names + copy",   "Low",    "Mechanical rename", AMBER, "YES"),
        ("Insight Rules",      "8 content rules",      "High",   "Per-vertical rules", RED, "YES"),
        ("Insight Rules",      "~40 generic rules",    "None",   "None",     GREEN,  "No"),
        ("Landing Page",       "Hero, features copy",  "Low",    "Per-vertical page", AMBER, "No"),
        ("Email Templates",    "Welcome tip copy",     "Low",    "Copy change", GREEN, "No"),
    ]

    col_x = [0.3, 2.45, 5.25, 7.0, 9.0, 11.0]
    col_w = [2.0, 2.7,  1.65, 1.85, 1.85, 2.0]
    headers = ["Category", "Component", "Effort", "Required Change", "Severity", "Blocks Launch?"]

    ry = 0.78
    for j, hdr in enumerate(headers):
        add_rect(slide, col_x[j], ry, col_w[j], 0.32, fill_color=RGBColor(0x27, 0x27, 0x2A))
        add_text(slide, hdr, col_x[j] + 0.08, ry + 0.04, col_w[j] - 0.1, 0.24,
                 size=8.5, bold=True, color=MUTED)

    sev_bg = {RED: RGBColor(0x3B, 0x10, 0x10), AMBER: RGBColor(0x2D, 0x1D, 0x00),
              GREEN: RGBColor(0x06, 0x25, 0x18), CYAN: RGBColor(0x06, 0x20, 0x28)}

    rh = 0.44
    for i, (cat, comp, effort, change, color, blocks) in enumerate(rows_data):
        ry2 = 1.14 + i * rh
        bg = sev_bg.get(color, CARD)
        add_rect(slide, col_x[0], ry2, sum(col_w) + 0.15, rh - 0.03, fill_color=bg)
        vals = [cat, comp, effort, change]
        for j, val in enumerate(vals):
            tc = WHITE if j <= 1 else MUTED
            add_text(slide, val, col_x[j] + 0.08, ry2 + 0.1, col_w[j] - 0.12, rh - 0.15,
                     size=9, color=tc)
        # severity dot
        add_rect(slide, col_x[4] + 0.2, ry2 + 0.09, 1.3, 0.25, fill_color=color)
        add_text(slide, "●", col_x[4] + 0.22, ry2 + 0.09, 0.28, 0.25,
                 size=8, color=WHITE)
        # blocks
        bc = RED if blocks == "YES" else MUTED
        btext = "YES" if blocks == "YES" else "No"
        add_text(slide, btext, col_x[5] + 0.5, ry2 + 0.1, 1.0, 0.25,
                 size=9, bold=True if blocks == "YES" else False, color=bc)


def slide_options_overview(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Verticalization Options — Overview",
                 "Four architecturally distinct approaches evaluated")

    options = [
        ("A", "Fully Separate\nCodebases",
         "Fork Prophet for each vertical. Independent Next.js apps, separate Supabase projects.",
         "Maximum isolation\nClean data separation\nSimple per-codebase",
         "Maintenance hell at 3+ verticals\nNo shared fixes\nDuplicate admin dashboards",
         RED, "NOT RECOMMENDED"),
        ("B", "Monorepo with\nVertical Packages",
         "Turborepo: @prophet/core shared packages + vertical-specific apps.",
         "Clean separation\nShared bug fixes\nIndependent deployments",
         "2–3 week infra project upfront\nOverkill for current team size",
         CYAN, "FUTURE — 3+ VERTICALS"),
        ("C", "Single Codebase +\nVertical Config Layer",
         "Single app, shared DB. VerticalConfig objects per vertical resolved at runtime.",
         "Fastest time to market\nOne deployment, one admin\nCross-vertical insights possible",
         "Discipline required for schema governance\nConfig layer can grow complex",
         GREEN, "RECOMMENDED"),
        ("D", "Single Codebase +\nSeparate Databases",
         "Same codebase, but each vertical gets its own Supabase project.",
         "Clean data separation\nVertical can be sold cleanly",
         "2× ops overhead\nMigrations applied N times\nHigher cost",
         AMBER, "NOT NOW"),
    ]

    opt_w = (W - 0.5 - 0.4 - 3 * 0.2) / 4
    for i, (letter, title, desc, pros, cons, color, verdict) in enumerate(options):
        cx = 0.4 + i * (opt_w + 0.2)
        oy = 0.75

        # Main card
        add_rect(slide, cx, oy, opt_w, 6.3, fill_color=CARD, border_color=color, border_width=1.0)
        add_rect(slide, cx, oy, opt_w, 0.07, fill_color=color)

        # Option letter badge
        add_rect(slide, cx + 0.15, oy + 0.12, 0.42, 0.42, fill_color=color)
        add_text(slide, letter, cx + 0.15, oy + 0.13, 0.42, 0.38,
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        add_text(slide, title, cx + 0.65, oy + 0.14, opt_w - 0.8, 0.55,
                 size=10, bold=True, color=WHITE)
        add_text(slide, desc,  cx + 0.12, oy + 0.72, opt_w - 0.25, 0.85,
                 size=9, color=MUTED)

        add_rect(slide, cx + 0.12, oy + 1.62, opt_w - 0.25, 0.025, fill_color=BORDER)

        add_text(slide, "PROS", cx + 0.12, oy + 1.68, opt_w - 0.25, 0.22,
                 size=8, bold=True, color=GREEN)
        for j, pro in enumerate(pros.split("\n")):
            add_text(slide, f"+ {pro}", cx + 0.12, oy + 1.9 + j * 0.3, opt_w - 0.25, 0.27,
                     size=8.5, color=GREEN)

        con_start = oy + 1.9 + len(pros.split("\n")) * 0.3 + 0.15
        add_text(slide, "CONS", cx + 0.12, con_start, opt_w - 0.25, 0.22,
                 size=8, bold=True, color=RED)
        for j, con in enumerate(cons.split("\n")):
            add_text(slide, f"− {con}", cx + 0.12, con_start + 0.22 + j * 0.3, opt_w - 0.25, 0.28,
                     size=8.5, color=MUTED)

        # Verdict badge at bottom
        add_rect(slide, cx + 0.12, oy + 5.72, opt_w - 0.25, 0.38, fill_color=color)
        add_text(slide, verdict, cx + 0.12, oy + 5.75, opt_w - 0.25, 0.32,
                 size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide_recommendation(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Recommendation: Option C",
                 "Single Codebase · Shared Database · Vertical Config Layer")

    badge(slide, "RECOMMENDED", W - 2.5, 0.18, w=2.2, h=0.3, bg=GREEN, size=9)

    phases = [
        ("Phase 1", "Schema Foundation",          "~1 day",  GREEN,
         "Add industry_type column to organizations\nBackfill existing orgs with 'restaurant'\nAdd index for query performance"),
        ("Phase 2", "Vertical Config System",      "~2–3 days", VIOLET,
         "Define VerticalConfig TypeScript interface\nExtract restaurant constants → lib/verticals/restaurant.ts\nWrite lib/verticals/liquor-store.ts\nWire getVerticalConfig() into all consumers"),
        ("Phase 3", "Onboarding Generalization",   "~1 day",  CYAN,
         "Rename restaurantName → businessName\nReplace CUISINES[] with verticalConfig.businessCategories\nReplace hardcoded copy with verticalConfig labels"),
        ("Phase 4", "Subdomain Routing",           "~1 day",  AMBER,
         "Create middleware.ts mapping hostnames → industry_type\nPass vertical through request context\nDNS: vaticliquor.com → same Vercel deployment"),
        ("Phase 5", "Vertical Landing Pages",      "~1 week", VIOLET_LT,
         "Liquor store landing page at /liquor or subdomain root\nParameterize animated SVG category labels\nExisting / remains restaurant landing"),
        ("Phase 6", "Liquor Store Content Signal", "~3–4 days", RED,
         "Liquor store product extraction Gemini prompt\nProduct category detection patterns\nLiquor promo keywords\n6 new catalog insight rules"),
    ]

    ph_w = 3.95
    ph_gap = 0.2
    ph_h = 3.0
    for i, (num, name, duration, color, body) in enumerate(phases):
        col = i % 3
        row = i // 3
        px = 0.3 + col * (ph_w + ph_gap)
        py = 0.75 + row * (ph_h + 0.2)
        add_rect(slide, px, py, ph_w, ph_h, fill_color=CARD, border_color=color, border_width=0.75)
        add_rect(slide, px, py, ph_w, 0.07, fill_color=color)
        add_text(slide, num, px + 0.15, py + 0.12, 0.8, 0.28, size=9, bold=True, color=color)
        add_text(slide, name, px + 0.15, py + 0.4, ph_w - 0.3, 0.35,
                 size=12, bold=True, color=WHITE)
        add_rect(slide, px + ph_w - 1.3, py + 0.1, 1.2, 0.28, fill_color=RGBColor(0x27, 0x27, 0x2A))
        add_text(slide, duration, px + ph_w - 1.3, py + 0.12, 1.2, 0.24,
                 size=8.5, color=MUTED, align=PP_ALIGN.CENTER)
        for j, line in enumerate(body.split("\n")):
            add_text(slide, f"→  {line}", px + 0.15, py + 0.85 + j * 0.42, ph_w - 0.25, 0.38,
                     size=9, color=MUTED)


def slide_what_stays(prs):
    slide = blank_slide(prs)
    slide_header(slide, "What Stays Untouched",
                 "The majority of the codebase requires zero changes — verticalization is additive, not destructive")

    badge(slide, "ZERO CHANGES", W - 2.4, 0.18, w=2.1, h=0.3, bg=GREEN, size=9)

    untouched = [
        ("Database Schema",         "All 15+ tables beyond the new industry_type column"),
        ("Billing & Stripe",        "Tier system, checkout, webhooks, trial period logic"),
        ("Auth & Sessions",         "Magic link, Google OAuth, Supabase auth callbacks"),
        ("Admin Dashboard",         "All 6 admin sections, waitlist, user/org management"),
        ("7 Intelligence Signals",  "Competitor, SEO, events, social, photos, traffic, weather"),
        ("~40+ Insight Rules",      "All non-content insight rules — social, SEO, events, traffic, weather, photos"),
        ("Background Job System",   "SSE streaming, ActiveJobBar, cron orchestrator"),
        ("Server-Side Caching",     "use cache directive, cacheTag/cacheLife, automatic revalidation"),
        ("Email Templates",         "5 of 6 templates are generic — only welcome tip copy changes"),
        ("Multi-Org Support",       "Org switcher, new org wizard, role system"),
        ("Insight Card System",     "Kanban view, status workflow, optimistic updates"),
        ("Supabase RLS Policies",   "All row-level security policies are org-scoped, not vertical-scoped"),
    ]

    item_w = 5.9
    item_h = 0.55
    gap_x = 0.4
    gap_y = 0.1
    per_row = 2

    for i, (name, desc) in enumerate(untouched):
        col = i % per_row
        row = i // per_row
        ix = 0.4 + col * (item_w + gap_x)
        iy = 0.78 + row * (item_h + gap_y)
        add_rect(slide, ix, iy, item_w, item_h,
                 fill_color=RGBColor(0x06, 0x1A, 0x12), border_color=GREEN, border_width=0.5)
        add_text(slide, "✓", ix + 0.1, iy + 0.12, 0.3, 0.3, size=13, bold=True, color=GREEN)
        add_text(slide, name, ix + 0.42, iy + 0.07, 2.2, 0.28, size=10, bold=True, color=WHITE)
        add_text(slide, desc, ix + 0.42, iy + 0.33, item_w - 0.55, 0.2, size=8.5, color=MUTED)


def slide_liquor_store(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Liquor Store Vertical — What Must Be Created",
                 "New content specific to the alcohol retail vertical — beyond cosmetic renaming")

    add_text(slide, "New Gemini Prompt — Product Catalog Extraction", 0.4, 0.75, 8, 0.28,
             size=12, bold=True, color=VIOLET_LT)

    prompt_items = [
        "Spirit type: bourbon, scotch, tequila, vodka, gin, rum, brandy",
        "Brand and distillery  ·  Price per bottle (750ml, 1L, 1.75L)  ·  ABV / proof",
        "Age statement (whiskeys)  ·  Country of origin  ·  Promotional vs. regular pricing",
    ]
    bullet_list(slide, prompt_items, 0.4, 1.05, 9, item_size=10, spacing=0.28)

    add_rect(slide, 0.4, 1.73, 12.5, 0.025, fill_color=BORDER)

    # Three columns: features, insight rules, new signals
    cols = [
        ("New Detected Features", CYAN, [
            "curbsidePickup",
            "homeDelivery",
            "loyaltyProgram",
            "tastingEvents",
            "privateLabelProducts",
            "bulkOrdering",
            "Drizly / Instacart integration",
        ]),
        ("New Catalog Insight Rules", VIOLET_LT, [
            "catalog.price_positioning_shift",
            "catalog.product_category_gap",
            "catalog.promo_signal_detected",
            "catalog.delivery_platform_gap",
            "catalog.exclusive_product_detected",
            "catalog.pricing_tier_gap  (new!)",
        ]),
        ("New Signals Worth Considering", AMBER, [
            "Distributor pricing intelligence",
            "Liquor license tracking (public record)",
            "State compliance / regulatory alerts",
            "Holiday sales law changes by state",
            "(Phase 2 — not MVP)",
        ]),
    ]

    col_w = 3.9
    for i, (title, color, items) in enumerate(cols):
        cx = 0.4 + i * (col_w + 0.28)
        add_rect(slide, cx, 1.83, col_w, 4.8, fill_color=CARD, border_color=color, border_width=0.75)
        add_rect(slide, cx, 1.83, col_w, 0.06, fill_color=color)
        add_text(slide, title, cx + 0.15, 1.9, col_w - 0.3, 0.3, size=10, bold=True, color=WHITE)
        for j, item in enumerate(items):
            add_text(slide, f"→  {item}", cx + 0.15, 2.28 + j * 0.48, col_w - 0.3, 0.4,
                     size=9.5, color=MUTED)

    add_rect(slide, 0.4, 6.75, 12.5, 0.025, fill_color=BORDER)
    add_text(slide,
             "Signals that require ZERO changes for liquor stores:  Competitor monitoring  ·  SEO  ·  Events  "
             "·  Social  ·  Photos  ·  Foot Traffic  ·  Weather",
             0.4, 6.82, 12.2, 0.35, size=9.5, color=GREEN, italic=True)


def slide_open_questions(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Open Questions",
                 "Decisions the team must make before implementation begins")

    questions = [
        ("Q1", "Who owns the liquor store build?",
         "Anand full build or Henry takes part of it?", VIOLET),
        ("Q2", "Branching strategy?",
         "feature-anand branch or new feature-verticalization branch?", CYAN),
        ("Q3", "How does industry_type get set at sign-up?",
         "Subdomain detection? Admin sets on waitlist approval? User selects in onboarding? All three?", AMBER),
        ("Q4", "Minimum viable signals at launch?",
         "Should liquor store launch without Content Intelligence (fastest), then add in sprint 2?", AMBER),
        ("Q5", "Waitlist for a new vertical?",
         "Same waitlist form at vaticliquor.com? Separate ClickUp project for tickets?", VIOLET_LT),
        ("Q6", "Pricing strategy for liquor stores?",
         "Same tiers as restaurants, or different limits given different data volumes?", CYAN),
        ("Q7", "Data separation comfort level?",
         "Is legal/finance OK with industry_type-filtered export, or do they require full DB isolation?", RED),
        ("Q8", "Fixed enum vs. open string for industry_type?",
         "Enum is safe & type-checked but needs migration per vertical. Open string + validation table is more flexible.", MUTED),
    ]

    q_w = 5.9
    q_h = 1.45
    for i, (num, question, detail, color) in enumerate(questions):
        col = i % 2
        row = i // 2
        qx = 0.4 + col * (q_w + 0.5)
        qy = 0.78 + row * (q_h + 0.15)
        add_rect(slide, qx, qy, q_w, q_h, fill_color=CARD, border_color=color, border_width=0.75)
        add_rect(slide, qx, qy, 0.65, q_h, fill_color=RGBColor(0x18, 0x18, 0x1F))
        add_text(slide, num, qx + 0.05, qy + 0.48, 0.56, 0.4,
                 size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, question, qx + 0.78, qy + 0.12, q_w - 0.9, 0.35,
                 size=11, bold=True, color=WHITE)
        add_text(slide, detail, qx + 0.78, qy + 0.52, q_w - 0.9, 0.8,
                 size=9.5, color=MUTED, wrap=True)


def slide_closing(prs):
    slide = blank_slide(prs)

    # Left accent
    add_rect(slide, 0, 0, 0.06, H, fill_color=VIOLET)

    # Faint background glow
    add_rect(slide, W - 4, 1.5, 4, 4, fill_color=RGBColor(0x10, 0x08, 0x20))

    add_text(slide, "VATIC", 0.5, 1.5, 4, 0.8, size=48, bold=True, color=VIOLET)
    add_text(slide, "Verticalization PRD — Summary", 0.5, 2.3, 9, 0.45,
             size=20, bold=True, color=WHITE)

    add_rect(slide, 0.5, 2.9, 8, 0.025, fill_color=BORDER)

    summary = [
        "The database schema is already ~90% vertical-agnostic.",
        "Only the Content/Menu signal needs a full per-vertical rewrite.",
        "Option C (single codebase + VerticalConfig layer) is the recommended path.",
        "6 phases from schema migration to live liquor store vertical.",
        "~40+ insight rules, all 7 non-content signals, billing, auth, and admin require no changes.",
    ]
    for i, line in enumerate(summary):
        add_text(slide, f"→  {line}", 0.5, 3.1 + i * 0.5, 10.5, 0.4,
                 size=12, color=VIOLET_LT if i == 2 else MUTED)

    add_text(slide, "Next Step: Present options to team · Decide on Q1–Q8 open questions · Begin Phase 1 schema migration",
             0.5, 6.0, 11, 0.35, size=11, bold=True, color=WHITE)

    add_text(slide, "April 7, 2026  ·  Anand Iyer  ·  Alive Labs / Prophet",
             0.5, H - 0.55, 10, 0.3, size=10, color=MUTED)


# ─── Main ────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    print("Building slides...")
    slide_title(prs)                ; print("  01/14  Title")
    slide_context(prs)              ; print("  02/14  Context & Background")
    slide_guiding_principles(prs)   ; print("  03/14  Guiding Principles")
    slide_audit_overview(prs)       ; print("  04/14  Audit Overview")
    slide_db_schema(prs)            ; print("  05/14  Database Schema")
    slide_type_system(prs)          ; print("  06/14  Type System")
    slide_signals(prs)              ; print("  07/14  Intelligence Signals")
    slide_prompts(prs)              ; print("  08/14  AI Prompts")
    slide_onboarding(prs)           ; print("  09/14  Onboarding Flow")
    slide_insight_rules(prs)        ; print("  10/14  Insight Rules")
    slide_impact_matrix(prs)        ; print("  11/14  Impact Matrix")
    slide_options_overview(prs)     ; print("  12/14  Options Overview")
    slide_recommendation(prs)       ; print("  13/14  Recommendation + Phases")
    slide_what_stays(prs)           ; print("  14/14  What Stays Untouched")
    slide_liquor_store(prs)         ; print("  15/16  Liquor Store Specifics")
    slide_open_questions(prs)       ; print("  16/17  Open Questions")
    slide_closing(prs)              ; print("  17/17  Closing")

    out = "/Users/anandiyer/Downloads/Vatic_Verticalization_PRD.pptx"
    prs.save(out)
    print(f"\nSaved → {out}")


if __name__ == "__main__":
    main()
